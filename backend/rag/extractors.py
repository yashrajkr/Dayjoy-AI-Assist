"""
Text extractors for the RAG pipeline.

Supported formats (best-effort — falls back to "" on failure):

  - PDF            → pypdf (pure Python, no system deps)
  - DOCX           → python-docx
  - PPTX           → python-pptx (optional)
  - XLSX/CSV       → openpyxl / csv
  - TXT/MD         → plain read
  - JSON           → stringified

Each extractor returns:
  - text: full extracted plain text
  - pages: list of per-page text (for PDF; for other formats, a single
           element with the full text)
  - sections: list of (section_title, text) tuples detected heuristically
  - metadata: dict of file-level metadata (title, author, page_count, etc.)

Optional dependencies are imported lazily so the module loads cleanly
even when they're not installed.
"""

from __future__ import annotations

import csv
import io
import json
import os
import re
from typing import Any, Dict, List, Optional, Tuple

# MIME types we know how to handle. Used by the upload endpoint to reject
# unsupported files early.
SUPPORTED_MIME_TYPES: Dict[str, str] = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xls",
    "application/msword": "doc",
    "application/vnd.ms-powerpoint": "ppt",
    "text/plain": "txt",
    "text/markdown": "md",
    "text/csv": "csv",
    "application/json": "json",
}

SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",
    ".pptx": "pptx",
    ".ppt": "ppt",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".txt": "txt",
    ".md": "md",
    ".markdown": "md",
    ".csv": "csv",
    ".json": "json",
}


def detect_kind(filename: str, mime_type: Optional[str]) -> str:
    """Pick the extractor kind from mime or filename extension."""
    if mime_type and mime_type in SUPPORTED_MIME_TYPES:
        return SUPPORTED_MIME_TYPES[mime_type]
    ext = os.path.splitext(filename or "")[1].lower()
    return SUPPORTED_EXTENSIONS.get(ext, "txt")


def extract_text(
    content: bytes,
    filename: str,
    mime_type: Optional[str] = None,
) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    """Extract text from `content` (raw bytes).

    Returns (full_text, pages, sections, metadata).
    `sections` is a list of (heading, body_text) tuples — used by the
    chunker to produce section-aware chunks.
    """
    kind = detect_kind(filename, mime_type)
    if kind == "pdf":
        return _extract_pdf(content)
    if kind == "docx":
        return _extract_docx(content)
    if kind == "pptx":
        return _extract_pptx(content)
    if kind == "xlsx":
        return _extract_xlsx(content)
    if kind == "csv":
        return _extract_csv(content)
    if kind == "json":
        return _extract_json(content)
    # txt / md / unknown → plain text
    return _extract_text(content)


def extract_metadata(content: bytes, filename: str, mime_type: Optional[str]) -> Dict[str, Any]:
    """Return only the metadata dict for the file (no full extraction)."""
    try:
        _, _, _, meta = extract_text(content, filename, mime_type)
        return meta
    except Exception as e:  # pragma: no cover
        return {"extractor_error": str(e)}


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError(
            "pypdf is required for PDF extraction. Install with: pip install pypdf"
        ) from e
    reader = PdfReader(io.BytesIO(content))
    pages: List[str] = []
    for page in reader.pages:
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        pages.append(txt)
    full = "\n\n".join(pages)
    meta = {
        "page_count": len(pages),
        "kind": "pdf",
        "title": (reader.metadata.title if reader.metadata else None),
        "author": (reader.metadata.author if reader.metadata else None),
    }
    sections = _split_into_sections(full)
    return full, pages, sections, meta


def _extract_docx(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as e:
        raise RuntimeError(
            "python-docx is required for DOCX extraction. Install with: pip install python-docx"
        ) from e
    doc = Document(io.BytesIO(content))
    paragraphs: List[str] = []
    sections: List[Tuple[str, str]] = []
    current_heading = "Introduction"
    current_body: List[str] = []

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        # Heuristic: Word "Heading 1" / "Heading 2" styles
        style = (para.style.name or "").lower() if para.style else ""
        if "heading" in style or text.isupper() and len(text) < 80:
            if current_body:
                sections.append((current_heading, "\n".join(current_body)))
            current_heading = text
            current_body = []
        else:
            paragraphs.append(text)
            current_body.append(text)
    if current_body:
        sections.append((current_heading, "\n".join(current_body)))

    full = "\n\n".join(paragraphs) or "\n\n".join(b for _, b in sections)
    if not sections:
        sections = _split_into_sections(full)
    return full, [full], sections, {"kind": "docx", "paragraph_count": len(paragraphs)}


def _extract_pptx(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for PPTX extraction. Install with: pip install python-pptx"
        ) from e
    prs = Presentation(io.BytesIO(content))
    pages: List[str] = []
    sections: List[Tuple[str, str]] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: List[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                texts.append(line)
        slide_text = "\n".join(texts)
        pages.append(slide_text)
        title = texts[0] if texts else f"Slide {i}"
        sections.append((title, slide_text))
    full = "\n\n".join(pages)
    return full, pages, sections, {"kind": "pptx", "slide_count": len(pages)}


def _extract_xlsx(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise RuntimeError(
            "openpyxl is required for XLSX extraction. Install with: pip install openpyxl"
        ) from e
    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    pages: List[str] = []
    sections: List[Tuple[str, str]] = []
    for sheet in wb.worksheets:
        rows: List[str] = []
        for row in sheet.iter_rows(values_only=True):
            row_str = ", ".join("" if v is None else str(v) for v in row)
            if row_str.strip(", "):
                rows.append(row_str)
        sheet_text = "\n".join(rows)
        pages.append(sheet_text)
        sections.append((sheet.title, sheet_text))
    full = "\n\n".join(pages)
    return full, pages, sections, {"kind": "xlsx", "sheet_count": len(pages)}


def _extract_csv(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows: List[str] = []
    for row in reader:
        row_str = ", ".join(row)
        if row_str.strip(", "):
            rows.append(row_str)
    full = "\n".join(rows)
    return full, [full], [("CSV", full)], {"kind": "csv", "row_count": len(rows)}


def _extract_json(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    text = content.decode("utf-8", errors="replace")
    try:
        data = json.loads(text)
        # Stringify nested objects into a readable form
        full = json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        full = text
    return full, [full], [("JSON", full)], {"kind": "json"}


def _extract_text(content: bytes) -> Tuple[str, List[str], List[Tuple[str, str]], Dict[str, Any]]:
    text = content.decode("utf-8", errors="replace")
    sections = _split_into_sections(text)
    return text, [text], sections, {"kind": "text"}


# ---------------------------------------------------------------------------
# Heuristic section detection
# ---------------------------------------------------------------------------

_HEADING_PATTERNS = [
    # Markdown ATX headings: # Title, ## Title
    re.compile(r"^\s{0,3}(#{1,6})\s+(.+?)\s*$", re.MULTILINE),
    # Markdown Setext headings (underline of = or -)
    re.compile(r"^\s*(.+?)\s*\n[=-]+\s*$", re.MULTILINE),
    # ALL CAPS short line (≤ 80 chars), at least 4 chars, no trailing period
    re.compile(r"^\s*([A-Z0-9][A-Z0-9 \-:&,/']{3,78})\s*$", re.MULTILINE),
    # Numbered section: "1. Title" or "1.1 Title"
    re.compile(r"^\s*(\d+(?:\.\d+){0,3})\.\s+([^\n]{3,80})\s*$", re.MULTILINE),
]


def _split_into_sections(text: str) -> List[Tuple[str, str]]:
    """Best-effort section splitter for plain text / markdown.

    Returns [(heading, body_text), ...]. When no headings are found, returns
    a single section with the full text.
    """
    if not text.strip():
        return [("Document", "")]

    # Collect all (start, heading) pairs from all patterns
    headings: List[Tuple[int, str]] = []
    for pat in _HEADING_PATTERNS:
        for m in pat.finditer(text):
            # Pick the last group (the actual title text) — fall back to
            # the whole match if no groups.
            title = (m.group(2) if m.lastindex and m.lastindex >= 2 else m.group(1)) or m.group(0)
            title = title.strip().rstrip(".")
            if title and len(title) <= 120:
                headings.append((m.start(), title))
    headings.sort()

    if not headings:
        return [("Document", text)]

    sections: List[Tuple[str, str]] = []
    for i, (start, title) in enumerate(headings):
        end = headings[i + 1][0] if i + 1 < len(headings) else len(text)
        body = text[start:end].strip()
        # Strip the heading line from the body for cleaner chunks
        body_lines = body.split("\n", 1)
        if len(body_lines) > 1:
            body = body_lines[1].strip()
        if body:
            sections.append((title, body))
    # Anything before the first heading becomes an "Introduction" section
    if headings and headings[0][0] > 0:
        intro = text[: headings[0][0]].strip()
        if intro:
            sections.insert(0, ("Introduction", intro))
    return sections or [("Document", text)]
